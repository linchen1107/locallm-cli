"""
檔案操作核心工具模組
提供讀取、寫入、編輯檔案的核心功能
"""

import os
from pathlib import Path
from typing import Optional

# PDF 相關導入（可選依賴）
try:
    import fitz  # PyMuPDF
    HAS_PDF_SUPPORT = True
except ImportError:
    HAS_PDF_SUPPORT = False

# Office 文件相關導入（可選依賴）
try:
    import docx  # python-docx for Word files
    HAS_WORD_SUPPORT = True
except ImportError:
    HAS_WORD_SUPPORT = False

try:
    import openpyxl  # for Excel files
    HAS_EXCEL_SUPPORT = True
except ImportError:
    HAS_EXCEL_SUPPORT = False

try:
    from pptx import Presentation  # python-pptx for PowerPoint files
    HAS_POWERPOINT_SUPPORT = True
except ImportError:
    HAS_POWERPOINT_SUPPORT = False

# 新增文件格式支援
try:
    import csv
    import yaml
    import tomllib  # Python 3.11+
    HAS_CSV_SUPPORT = True
    HAS_YAML_SUPPORT = True
    HAS_TOML_SUPPORT = True
except ImportError:
    HAS_CSV_SUPPORT = True  # csv 是標準庫
    HAS_YAML_SUPPORT = False
    HAS_TOML_SUPPORT = False

# SQL 文件支援
try:
    import sqlite3
    HAS_SQL_SUPPORT = True
except ImportError:
    HAS_SQL_SUPPORT = False


class FileTools:
    """檔案操作工具類"""
    
    def __init__(self, base_path: Optional[str] = None):
        """
        初始化檔案工具
        
        Args:
            base_path: 基礎路徑，如果未提供則使用當前工作目錄
        """
        self.base_path = Path(base_path) if base_path else Path.cwd()
        
    def get_current_path(self) -> str:
        """取得目前的工作路徑"""
        return str(self.base_path.absolute())
    
    def _resolve_path(self, file_path: str) -> Path:
        """
        解析檔案路徑，處理相對路徑和絕對路徑
        
        Args:
            file_path: 檔案路徑
            
        Returns:
            Path: 解析後的完整路徑
        """
        path = Path(file_path)
        if path.is_absolute():
            return path
        else:
            return self.base_path / path
    
    def read_file(self, file_path: str) -> str:
        """
        讀取指定路徑的檔案內容
        
        Args:
            file_path: 檔案路徑
            
        Returns:
            str: 檔案內容
            
        Raises:
            FileNotFoundError: 檔案不存在
            PermissionError: 沒有讀取權限
            UnicodeDecodeError: 檔案編碼問題
        """
        resolved_path = self._resolve_path(file_path)
        
        if not resolved_path.exists():
            raise FileNotFoundError(f"檔案不存在: {resolved_path}")
        
        if not resolved_path.is_file():
            raise ValueError(f"路徑不是檔案: {resolved_path}")
        
        try:
            # 嘗試使用 UTF-8 編碼讀取
            with open(resolved_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return content
        except UnicodeDecodeError:
            try:
                # 如果 UTF-8 失敗，嘗試使用系統預設編碼
                with open(resolved_path, 'r', encoding='utf-8-sig') as f:
                    content = f.read()
                return content
            except UnicodeDecodeError:
                # 最後嘗試使用 latin-1 編碼（幾乎不會失敗）
                with open(resolved_path, 'r', encoding='latin-1') as f:
                    content = f.read()
                return content

    def read_pdf(self, file_path: str, use_ocr: bool = False, extract_images: bool = False) -> str:
        """
        讀取 PDF 檔案內容並提取文字，特別優化論文處理
        
        Args:
            file_path: PDF 檔案路徑
            use_ocr: 是否強制使用 OCR（適用於掃描型 PDF）
            extract_images: 是否提取圖片信息
            
        Returns:
            str: 提取的文字內容
            
        Raises:
            FileNotFoundError: 檔案不存在
            ValueError: 不是 PDF 檔案或 PDF 讀取失敗
            ImportError: PyMuPDF 未安裝
        """
        if not HAS_PDF_SUPPORT:
            raise ImportError("PDF 支援需要安裝 PyMuPDF。請執行: pip install pymupdf")
        
        resolved_path = self._resolve_path(file_path)
        
        if not resolved_path.exists():
            raise FileNotFoundError(f"檔案不存在: {resolved_path}")
        
        if not resolved_path.is_file():
            raise ValueError(f"路徑不是檔案: {resolved_path}")
        
        if resolved_path.suffix.lower() != '.pdf':
            raise ValueError(f"檔案不是 PDF 格式: {resolved_path}")
        
        try:
            # 使用 PyMuPDF 讀取 PDF
            pdf_document = fitz.open(str(resolved_path))
            text_content = ""
            low_text_pages = []  # 記錄文字內容較少的頁面
            image_count = 0
            math_formula_count = 0
            
            # 逐頁提取文字
            for page_num in range(pdf_document.page_count):
                page = pdf_document[page_num]
                page_header = f"--- 第 {page_num + 1} 頁 ---\n"
                
                # 使用 get_text() 方法提取文字
                page_text = getattr(page, 'get_text', lambda: "")()
                
                # 檢查文字內容是否充足
                if len(page_text.strip()) < 50:  # 文字內容較少，可能需要 OCR
                    low_text_pages.append(page_num)
                
                if use_ocr or (len(page_text.strip()) < 50):
                    # 嘗試使用 OCR
                    ocr_text = self._try_ocr_on_page(str(resolved_path), page_num)
                    if ocr_text and len(ocr_text.strip()) > len(page_text.strip()):
                        page_text = ocr_text
                        page_header += "[OCR] "
                
                # 檢測數學公式（簡單檢測）
                if any(formula in page_text for formula in ['∑', '∫', '∂', '∇', '∞', 'α', 'β', 'γ', 'δ', 'ε', 'ζ', 'η', 'θ', 'λ', 'μ', 'π', 'ρ', 'σ', 'τ', 'φ', 'χ', 'ψ', 'ω']):
                    math_formula_count += 1
                    page_header += "[含數學公式] "
                
                # 提取圖片信息
                if extract_images:
                    image_list = page.get_images()
                    if image_list:
                        image_count += len(image_list)
                        page_header += f"[含 {len(image_list)} 張圖片] "
                
                text_content += page_header + page_text + "\n\n"
            
            # 保存頁數信息
            total_pages = pdf_document.page_count
            pdf_document.close()
            
            # 添加文檔摘要信息
            summary_info = f"\n\n=== PDF 文檔摘要 ===\n"
            summary_info += f"總頁數: {total_pages}\n"
            summary_info += f"總文字長度: {len(text_content)} 字符\n"
            if image_count > 0:
                summary_info += f"圖片數量: {image_count}\n"
            if math_formula_count > 0:
                summary_info += f"含數學公式頁數: {math_formula_count}\n"
            if low_text_pages:
                summary_info += f"低文字內容頁數: {len(low_text_pages)} (頁碼: {', '.join(str(p+1) for p in low_text_pages)})\n"
            
            # 判斷文檔類型
            if any(keyword in text_content.lower() for keyword in ['abstract', '摘要', 'introduction', '引言', 'conclusion', '結論', 'references', '參考文獻']):
                summary_info += "文檔類型: 學術論文\n"
            elif any(keyword in text_content.lower() for keyword in ['chapter', '章節', 'section', '節']):
                summary_info += "文檔類型: 書籍/報告\n"
            else:
                summary_info += "文檔類型: 一般文檔\n"
            
            text_content += summary_info
            
            # 如果整體文字內容很少，提示可能需要 OCR
            if not text_content.strip() or len(text_content.strip()) < 100:
                if low_text_pages:
                    ocr_hint = f"\n\n（注意：第 {', '.join(str(p+1) for p in low_text_pages)} 頁文字內容較少，"
                    ocr_hint += "可能是掃描型 PDF，建議使用 OCR 功能提取文字）"
                    text_content += ocr_hint
                else:
                    text_content = "（此 PDF 文件沒有可提取的文字內容，可能是圖片格式的 PDF，建議使用 OCR 功能）"
            
            return text_content.strip()
            
        except Exception as e:
            raise ValueError(f"PDF 讀取失敗: {e}")
    
    def _try_ocr_on_page(self, pdf_path: str, page_num: int) -> str:
        """
        嘗試對 PDF 頁面使用 OCR
        
        Args:
            pdf_path: PDF 檔案路徑
            page_num: 頁面編號
            
        Returns:
            str: OCR 提取的文字，失敗則返回空字串
        """
        try:
            from .ocr_tools import create_ocr_processor
            
            ocr_processor = create_ocr_processor()
            if ocr_processor:
                return ocr_processor.extract_text_from_pdf_page(pdf_path, page_num)
        except ImportError:
            pass  # OCR 不可用，靜默忽略
        except Exception:
            pass  # OCR 失敗，靜默忽略
        
        return ""
    
    def read_word(self, file_path: str) -> str:
        """
        讀取 Word 檔案內容並提取文字
        
        Args:
            file_path: Word 檔案路徑 (.docx)
            
        Returns:
            str: 提取的文字內容
            
        Raises:
            FileNotFoundError: 檔案不存在
            ValueError: 不是 Word 檔案或讀取失敗
            ImportError: python-docx 未安裝
        """
        if not HAS_WORD_SUPPORT:
            raise ImportError("Word 支援需要安裝 python-docx。請執行: pip install python-docx")
        
        resolved_path = self._resolve_path(file_path)
        
        if not resolved_path.exists():
            raise FileNotFoundError(f"檔案不存在: {resolved_path}")
        
        if not resolved_path.is_file():
            raise ValueError(f"路徑不是檔案: {resolved_path}")
        
        if resolved_path.suffix.lower() not in ['.docx']:
            raise ValueError(f"檔案不是支援的 Word 格式 (.docx): {resolved_path}")
        
        try:
            doc = docx.Document(str(resolved_path))
            text_content = ""
            
            # 提取段落文字
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            # 提取表格文字
            for table in doc.tables:
                text_content += "\n--- 表格 ---\n"
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    text_content += " | ".join(row_text) + "\n"
                text_content += "\n"
            
            return text_content.strip()
        
        except Exception as e:
            raise ValueError(f"Word 檔案讀取失敗: {e}")
    
    def read_excel(self, file_path: str, sheet_name: Optional[str] = None) -> str:
        """
        讀取 Excel 檔案內容並提取文字
        
        Args:
            file_path: Excel 檔案路徑 (.xlsx, .xls)
            sheet_name: 工作表名稱，如果未指定則讀取第一個工作表
            
        Returns:
            str: 提取的文字內容
            
        Raises:
            FileNotFoundError: 檔案不存在
            ValueError: 不是 Excel 檔案或讀取失敗
            ImportError: openpyxl 未安裝
        """
        if not HAS_EXCEL_SUPPORT:
            raise ImportError("Excel 支援需要安裝 openpyxl。請執行: pip install openpyxl")
        
        resolved_path = self._resolve_path(file_path)
        
        if not resolved_path.exists():
            raise FileNotFoundError(f"檔案不存在: {resolved_path}")
        
        if not resolved_path.is_file():
            raise ValueError(f"路徑不是檔案: {resolved_path}")
        
        if resolved_path.suffix.lower() not in ['.xlsx', '.xlsm']:
            raise ValueError(f"檔案不是支援的 Excel 格式 (.xlsx, .xlsm): {resolved_path}")
        
        try:
            workbook = openpyxl.load_workbook(str(resolved_path), data_only=True)
            text_content = ""
            
            # 選擇工作表
            if sheet_name:
                if sheet_name not in workbook.sheetnames:
                    raise ValueError(f"工作表 '{sheet_name}' 不存在。可用工作表: {workbook.sheetnames}")
                worksheet = workbook[sheet_name]
                text_content += f"=== 工作表：{sheet_name} ===\n\n"
            else:
                worksheet = workbook.active
                if worksheet and hasattr(worksheet, 'title'):
                    text_content += f"=== 工作表：{worksheet.title} ===\n\n"
                else:
                    text_content += f"=== 預設工作表 ===\n\n"
            
            # 提取儲存格內容
            if worksheet and hasattr(worksheet, 'iter_rows'):
                for row in worksheet.iter_rows(values_only=True):
                    row_text = []
                    for cell_value in row:
                        if cell_value is not None:
                            row_text.append(str(cell_value))
                        else:
                            row_text.append("")
                    
                    # 只添加非空行
                    if any(cell.strip() for cell in row_text if cell):
                        text_content += " | ".join(row_text) + "\n"
            
            # 如果有多個工作表，顯示工作表列表
            if len(workbook.sheetnames) > 1 and worksheet:
                worksheet_title = getattr(worksheet, 'title', '')
                other_sheets = [name for name in workbook.sheetnames if name != worksheet_title]
                if other_sheets:
                    text_content += f"\n--- 其他工作表：{', '.join(other_sheets)} ---\n"
            
            workbook.close()
            return text_content.strip()
        
        except Exception as e:
            raise ValueError(f"Excel 檔案讀取失敗: {e}")
    
    def read_powerpoint(self, file_path: str) -> str:
        """
        讀取 PowerPoint 檔案內容並提取文字
        
        Args:
            file_path: PowerPoint 檔案路徑 (.pptx)
            
        Returns:
            str: 提取的文字內容
            
        Raises:
            FileNotFoundError: 檔案不存在
            ValueError: 不是 PowerPoint 檔案或讀取失敗
            ImportError: python-pptx 未安裝
        """
        if not HAS_POWERPOINT_SUPPORT:
            raise ImportError("PowerPoint 支援需要安裝 python-pptx。請執行: pip install python-pptx")
        
        resolved_path = self._resolve_path(file_path)
        
        if not resolved_path.exists():
            raise FileNotFoundError(f"檔案不存在: {resolved_path}")
        
        if not resolved_path.is_file():
            raise ValueError(f"路徑不是檔案: {resolved_path}")
        
        if resolved_path.suffix.lower() not in ['.pptx']:
            raise ValueError(f"檔案不是支援的 PowerPoint 格式 (.pptx): {resolved_path}")
        
        try:
            presentation = Presentation(str(resolved_path))
            text_content = ""
            
            # 逐張投影片提取文字
            for slide_num, slide in enumerate(presentation.slides, 1):
                text_content += f"=== 投影片 {slide_num} ===\n"
                
                # 提取投影片中的文字
                for shape in slide.shapes:
                    # 安全地檢查和提取文字內容
                    try:
                        if hasattr(shape, 'text_frame'):
                            text_frame = getattr(shape, 'text_frame', None)
                            if text_frame:
                                shape_text = ""
                                for paragraph in text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        shape_text += run.text
                                    shape_text += "\n"
                                if shape_text.strip():
                                    text_content += shape_text
                        elif hasattr(shape, 'text'):
                            shape_text_attr = getattr(shape, 'text', None)
                            if shape_text_attr:
                                shape_text = str(shape_text_attr)
                                if shape_text.strip():
                                    text_content += shape_text + "\n"
                    except (AttributeError, TypeError, Exception):
                        pass  # 忽略無法訪問文字的 shape
                    
                    # 如果是表格，提取表格內容
                    try:
                        if hasattr(shape, 'table'):
                            table = getattr(shape, 'table', None)
                            if table:
                                text_content += "\n--- 表格 ---\n"
                                for row in table.rows:
                                    row_text = []
                                    for cell in row.cells:
                                        cell_text_attr = getattr(cell, 'text', "")
                                        cell_text = str(cell_text_attr) if cell_text_attr else ""
                                        row_text.append(cell_text.strip())
                                    text_content += " | ".join(row_text) + "\n"
                                text_content += "\n"
                    except (AttributeError, TypeError, Exception):
                        pass  # 忽略無法訪問表格的 shape
                
                text_content += "\n"
            
            return text_content.strip()
        
        except Exception as e:
            raise ValueError(f"PowerPoint 檔案讀取失敗: {e}")
    
    def write_file(self, file_path: str, content: str) -> None:
        """
        將內容寫入指定檔案（覆蓋模式）
        
        Args:
            file_path: 檔案路徑
            content: 要寫入的內容
            
        Raises:
            PermissionError: 沒有寫入權限
            OSError: 其他系統錯誤
        """
        resolved_path = self._resolve_path(file_path)
        
        # 確保父目錄存在
        resolved_path.parent.mkdir(parents=True, exist_ok=True)
        
        with open(resolved_path, 'w', encoding='utf-8') as f:
            f.write(content)
    
    def edit_file(self, file_path: str, new_content: str) -> None:
        """
        編輯檔案內容（覆蓋模式）
        這個函數在語義上與 write_file 相同，但用於明確表示是在編輯現有檔案
        
        Args:
            file_path: 檔案路徑
            new_content: 新的檔案內容
            
        Raises:
            PermissionError: 沒有寫入權限
            OSError: 其他系統錯誤
        """
        self.write_file(file_path, new_content)
    
    def file_exists(self, file_path: str) -> bool:
        """
        檢查檔案是否存在
        
        Args:
            file_path: 檔案路徑
            
        Returns:
            bool: 檔案是否存在
        """
        resolved_path = self._resolve_path(file_path)
        return resolved_path.exists() and resolved_path.is_file()
    
    def list_files(self, directory: str = ".", pattern: str = "*") -> list:
        """
        列出目錄中的檔案
        
        Args:
            directory: 目錄路徑，預設為當前目錄
            pattern: 檔案模式，預設為所有檔案
            
        Returns:
            list: 檔案路徑列表
        """
        resolved_path = self._resolve_path(directory)
        
        if not resolved_path.exists():
            raise FileNotFoundError(f"目錄不存在: {resolved_path}")
        
        if not resolved_path.is_dir():
            raise ValueError(f"路徑不是目錄: {resolved_path}")
        
        files = []
        for file_path in resolved_path.glob(pattern):
            if file_path.is_file():
                files.append(str(file_path.relative_to(self.base_path)))
        
        return sorted(files)
    
    def get_file_info(self, file_path: str) -> dict:
        """
        取得檔案資訊
        
        Args:
            file_path: 檔案路徑
            
        Returns:
            dict: 檔案資訊字典
        """
        resolved_path = self._resolve_path(file_path)
        
        if not resolved_path.exists():
            raise FileNotFoundError(f"檔案不存在: {resolved_path}")
        
        stat = resolved_path.stat()
        
        return {
            'path': str(resolved_path),
            'name': resolved_path.name,
            'size': stat.st_size,
            'modified': stat.st_mtime,
            'is_file': resolved_path.is_file(),
            'is_dir': resolved_path.is_dir(),
            'extension': resolved_path.suffix
        }
    
    def read_csv(self, file_path: str) -> str:
        """
        讀取 CSV 檔案內容並格式化顯示
        
        Args:
            file_path: CSV 檔案路徑
            
        Returns:
            str: 格式化的 CSV 內容
        """
        resolved_path = self._resolve_path(file_path)
        
        if not resolved_path.exists():
            raise FileNotFoundError(f"檔案不存在: {resolved_path}")
        
        if resolved_path.suffix.lower() != '.csv':
            raise ValueError(f"檔案不是 CSV 格式: {resolved_path}")
        
        try:
            content = ""
            with open(resolved_path, 'r', encoding='utf-8', newline='') as f:
                csv_reader = csv.reader(f)
                rows = list(csv_reader)
                
                if not rows:
                    return "CSV 檔案為空"
                
                # 格式化顯示
                max_widths = [max(len(str(cell)) for cell in col) for col in zip(*rows)]
                
                for i, row in enumerate(rows):
                    formatted_row = []
                    for j, cell in enumerate(row):
                        formatted_cell = str(cell).ljust(max_widths[j] if j < len(max_widths) else 0)
                        formatted_row.append(formatted_cell)
                    
                    content += " | ".join(formatted_row) + "\n"
                    
                    # 添加分隔線（僅在第一行後）
                    if i == 0:
                        separator = []
                        for width in max_widths:
                            separator.append("-" * width)
                        content += " | ".join(separator) + "\n"
            
            return content.strip()
            
        except Exception as e:
            raise ValueError(f"CSV 檔案讀取失敗: {e}")
    
    def read_yaml(self, file_path: str) -> str:
        """
        讀取 YAML 檔案內容
        
        Args:
            file_path: YAML 檔案路徑
            
        Returns:
            str: YAML 內容
        """
        if not HAS_YAML_SUPPORT:
            raise ImportError("YAML 支援需要安裝 PyYAML。請執行: pip install pyyaml")
        
        resolved_path = self._resolve_path(file_path)
        
        if not resolved_path.exists():
            raise FileNotFoundError(f"檔案不存在: {resolved_path}")
        
        if resolved_path.suffix.lower() not in ['.yml', '.yaml']:
            raise ValueError(f"檔案不是 YAML 格式: {resolved_path}")
        
        try:
            with open(resolved_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # 驗證 YAML 格式
            yaml.safe_load(content)
            return content
            
        except Exception as e:
            raise ValueError(f"YAML 檔案讀取失敗: {e}")
    
    def read_toml(self, file_path: str) -> str:
        """
        讀取 TOML 檔案內容
        
        Args:
            file_path: TOML 檔案路徑
            
        Returns:
            str: TOML 內容
        """
        if not HAS_TOML_SUPPORT:
            raise ImportError("TOML 支援需要 Python 3.11+ 或安裝 toml。請執行: pip install toml")
        
        resolved_path = self._resolve_path(file_path)
        
        if not resolved_path.exists():
            raise FileNotFoundError(f"檔案不存在: {resolved_path}")
        
        if resolved_path.suffix.lower() != '.toml':
            raise ValueError(f"檔案不是 TOML 格式: {resolved_path}")
        
        try:
            with open(resolved_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # 驗證 TOML 格式
            tomllib.loads(content)
            return content
            
        except Exception as e:
            raise ValueError(f"TOML 檔案讀取失敗: {e}")
    
    def read_sql(self, file_path: str) -> str:
        """
        讀取 SQL 檔案內容
        
        Args:
            file_path: SQL 檔案路徑
            
        Returns:
            str: SQL 內容
        """
        resolved_path = self._resolve_path(file_path)
        
        if not resolved_path.exists():
            raise FileNotFoundError(f"檔案不存在: {resolved_path}")
        
        if resolved_path.suffix.lower() != '.sql':
            raise ValueError(f"檔案不是 SQL 格式: {resolved_path}")
        
        try:
            with open(resolved_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            return content
            
        except Exception as e:
            raise ValueError(f"SQL 檔案讀取失敗: {e}")


# 創建預設實例
default_file_tools = FileTools()

# 提供便捷的函數介面
def read_file(file_path: str) -> str:
    """讀取檔案內容"""
    return default_file_tools.read_file(file_path)

def write_file(file_path: str, content: str) -> None:
    """寫入檔案內容"""
    default_file_tools.write_file(file_path, content)

def edit_file(file_path: str, new_content: str) -> None:
    """編輯檔案內容"""
    default_file_tools.edit_file(file_path, new_content)

def file_exists(file_path: str) -> bool:
    """檢查檔案是否存在"""
    return default_file_tools.file_exists(file_path)

def list_files(directory: str = ".", pattern: str = "*") -> list:
    """列出目錄中的檔案"""
    return default_file_tools.list_files(directory, pattern)

def get_current_path() -> str:
    """取得目前工作路徑"""
    return default_file_tools.get_current_path()

def read_pdf(file_path: str, use_ocr: bool = False, extract_images: bool = False) -> str:
    """讀取 PDF 檔案內容"""
    return default_file_tools.read_pdf(file_path, use_ocr, extract_images)

def read_csv(file_path: str) -> str:
    """讀取 CSV 檔案內容"""
    return default_file_tools.read_csv(file_path)

def read_yaml(file_path: str) -> str:
    """讀取 YAML 檔案內容"""
    return default_file_tools.read_yaml(file_path)

def read_toml(file_path: str) -> str:
    """讀取 TOML 檔案內容"""
    return default_file_tools.read_toml(file_path)

def read_sql(file_path: str) -> str:
    """讀取 SQL 檔案內容"""
    return default_file_tools.read_sql(file_path)